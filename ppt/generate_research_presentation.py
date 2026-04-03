#!/usr/bin/env python3
"""
Cloud Threat Mapper - Research Presentation (10 Slides)
Modern dark theme with professional design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Modern Dark Theme Color Palette
COLORS = {
    'bg_dark': RGBColor(13, 25, 45),         # Deep navy blue background
    'bg_card': RGBColor(22, 39, 68),         # Slightly lighter for cards
    'bg_card_hover': RGBColor(30, 52, 89),   # Hover state
    'primary': RGBColor(56, 189, 248),       # Bright cyan/blue
    'primary_dark': RGBColor(14, 165, 233),  # Darker blue
    'accent': RGBColor(139, 92, 246),        # Purple accent
    'success': RGBColor(34, 197, 94),        # Green
    'warning': RGBColor(251, 146, 60),       # Orange
    'danger': RGBColor(248, 113, 113),       # Red/Pink
    'text_white': RGBColor(255, 255, 255),   # Pure white text
    'text_light': RGBColor(199, 215, 232),   # Light gray-blue
    'text_muted': RGBColor(148, 163, 184),   # Muted text
    'border': RGBColor(51, 65, 85),          # Subtle borders
    'grid': RGBColor(30, 41, 59),            # Grid lines
}


def set_dark_background(slide):
    """Set dark navy background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg_dark']


def add_decorative_elements(slide):
    """Add subtle decorative geometric elements."""
    # Corner accents
    # Top-left
    for i in range(4):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3 + i * 0.15), Inches(0.3),
            Inches(0.12), Inches(0.12)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['primary']
        rect.fill.transparency = 0.6 - i * 0.1
        rect.line.fill.background()

    # Top-right
    for i in range(4):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9.5 - i * 0.15), Inches(0.3),
            Inches(0.12), Inches(0.12)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['primary']
        rect.fill.transparency = 0.6 - i * 0.1
        rect.line.fill.background()

    # Bottom-left
    for i in range(4):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3 + i * 0.15), Inches(7.0),
            Inches(0.12), Inches(0.12)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['accent']
        rect.fill.transparency = 0.6 - i * 0.1
        rect.line.fill.background()

    # Bottom-right
    for i in range(4):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9.5 - i * 0.15), Inches(7.0),
            Inches(0.12), Inches(0.12)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['accent']
        rect.fill.transparency = 0.6 - i * 0.1
        rect.line.fill.background()

    # Scattered decorative elements
    deco_positions = [
        (Inches(2), Inches(1), COLORS['primary'], 0.8),
        (Inches(7), Inches(1.5), COLORS['accent'], 0.7),
        (Inches(1.5), Inches(5), COLORS['accent'], 0.7),
        (Inches(8), Inches(5.5), COLORS['primary'], 0.8),
    ]

    for x, y, color, transparency in deco_positions:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(0.15), Inches(0.15)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = transparency
        circle.line.fill.background()


def add_title_slide(prs):
    """Slide 1: Modern title slide with dark theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Large background gradient effect (subtle circles)
    for i in range(3):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6 + i * 1.5), Inches(-1 + i * 0.5),
            Inches(3 - i * 0.5), Inches(3 - i * 0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary']
        circle.fill.transparency = 0.85 + i * 0.05
        circle.line.fill.background()

    # Left accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['primary']
    accent.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2), Inches(8.5), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cloud Threat Mapper"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.3), Inches(8.5), Inches(1)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Graph-Based Attack Path Analysis for"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    p2 = tf.add_paragraph()
    p2.text = "Cloud Infrastructure Security"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLORS['text_light']
    p2.font.name = 'Segoe UI'

    # Tagline
    tag_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.5), Inches(8.5), Inches(0.5)
    )
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Automated detection of privilege escalation paths using graph algorithms and AI"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_muted']
    p.font.name = 'Segoe UI'

    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(6.8), Inches(8.5), Inches(0.4)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Research Project Presentation | 2026"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_muted']
    p.font.name = 'Segoe UI'

    return slide


def add_problem_statement_slide(prs):
    """Slide 2: Problem Statement with clean cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "02"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Problem Statement"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Problem cards - vertically stacked
    problems = [
        (
            "Cloud Infrastructure Complexity",
            "Modern cloud environments have hundreds of interconnected resources with complex dependency chains",
            COLORS['danger']
        ),
        (
            "Privilege Escalation Detection Gap",
            "Traditional scanners miss multi-hop attack paths spanning IAM, compute, and storage",
            COLORS['warning']
        ),
        (
            "Rule-Based Scanner Limitations",
            "Checklist-based tools produce long reports without context or prioritization",
            COLORS['primary']
        ),
    ]

    for i, (title, desc, color) in enumerate(problems):
        y = Inches(1.3 + i * 1.75)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(8.6), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_card']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Left accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.15), Inches(1.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # Title
        ct = card.text_frame
        ct.paragraphs[0].text = title
        ct.paragraphs[0].font.size = Pt(15)
        ct.paragraphs[0].font.bold = True
        ct.paragraphs[0].font.color.rgb = color
        ct.paragraphs[0].space_left = Pt(50)

        # Description
        p = ct.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)
        p.space_left = Pt(50)

    return slide


def add_research_question_slide(prs):
    """Slide 3: Research Question."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "03"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Research Question"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Main question box
    q_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.3), Inches(8.6), Inches(1.8)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = COLORS['bg_card']
    q_box.line.color.rgb = COLORS['primary']
    q_box.line.width = Pt(2)

    qt = q_box.text_frame
    qt.paragraphs[0].text = "How can graph-based infrastructure modeling improve the detection of privilege escalation paths in cloud environments?"
    qt.paragraphs[0].font.size = Pt(18)
    qt.paragraphs[0].font.bold = True
    qt.paragraphs[0].font.color.rgb = COLORS['text_white']
    qt.paragraphs[0].alignment = PP_ALIGN.CENTER
    qt.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Supporting points
    points = [
        ("Cloud Resource Relationships", "Understanding how IAM, compute, and storage interconnect"),
        ("Attack Surface Expansion", "Mapping multi-hop privilege escalation chains"),
        ("Automated Reasoning", "AI-assisted analysis for risk prioritization"),
    ]

    for i, (point, desc) in enumerate(points):
        x = Inches(0.7 + i * 2.95)
        y = Inches(3.5)

        # Point card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_card']
        card.line.color.rgb = COLORS['primary']
        card.line.width = Pt(1.5)

        # Number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS['primary']
        badge.line.fill.background()
        bt = badge.text_frame
        bt.paragraphs[0].text = str(i + 1)
        bt.paragraphs[0].font.size = Pt(16)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].font.color.rgb = COLORS['bg_dark']
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Point text
        ct = card.text_frame
        ct.paragraphs[0].text = ""
        p = ct.add_paragraph()
        p.text = point
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        p.space_left = Pt(45)

        p2 = ct.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLORS['text_muted']
        p2.space_before = Pt(4)
        p2.space_left = Pt(45)

    return slide


def add_proposed_approach_slide(prs):
    """Slide 4: Proposed Approach with clean pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "04"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Proposed Approach"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Pipeline stages - horizontal
    stages = [
        ("1", "Infrastructure\nData Collection", "AWS APIs via boto3", COLORS['primary']),
        ("2", "Graph-Based\nModeling", "Resources as nodes", COLORS['accent']),
        ("3", "Attack Path\nDiscovery", "DFS/BFS traversal", COLORS['warning']),
        ("4", "AI-Assisted\nReasoning", "LLM risk analysis", COLORS['success']),
    ]

    for i, (num, title, desc, color) in enumerate(stages):
        x = Inches(0.7 + i * 2.2)
        y = Inches(1.5)

        # Connector line
        if i < len(stages) - 1:
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(2), y + Inches(1.25),
                Inches(0.15), Inches(0.08)
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = COLORS['text_muted']
            conn.line.fill.background()

        # Stage card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.05), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_card']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Number circle
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.725), y + Inches(0.15), Inches(0.6), Inches(0.6)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        bt = badge.text_frame
        bt.paragraphs[0].text = num
        bt.paragraphs[0].font.size = Pt(20)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].font.color.rgb = COLORS['bg_dark']
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        ct = card.text_frame
        ct.paragraphs[0].text = ""
        p = ct.add_paragraph()
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(45)

        # Description
        p2 = ct.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLORS['text_muted']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

    return slide


def add_system_architecture_slide(prs):
    """Slide 5: System Architecture with clean diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "05"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Architecture components - clean layered layout
    layers = [
        # Frontend layer
        (Inches(3.5), Inches(1.3), Inches(3), Inches(0.6),
         "Frontend UI", "React + Cytoscape.js", COLORS['primary']),

        # API layer
        (Inches(3.5), Inches(2.2), Inches(3), Inches(0.6),
         "API Layer", "FastAPI REST Server", COLORS['success']),

        # Worker layer
        (Inches(3.5), Inches(3.1), Inches(3), Inches(0.6),
         "Task Queue", "Celery + Redis", COLORS['warning']),

        # Data layer - 3 columns
        (Inches(0.7), Inches(4.1), Inches(2.7), Inches(1.1),
         "PostgreSQL", "Scan metadata\n& findings", COLORS['primary']),

        (Inches(3.65), Inches(4.1), Inches(2.7), Inches(1.1),
         "Neo4j", "Graph storage\nDFS/BFS traversal", COLORS['accent']),

        (Inches(6.6), Inches(4.1), Inches(2.7), Inches(1.1),
         "AI Engine", "LLM analysis", COLORS['success']),
    ]

    for x, y, w, h, title, desc, color in layers:
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_card']
        box.line.color.rgb = color
        box.line.width = Pt(2)

        bt = box.text_frame
        bt.paragraphs[0].text = title
        bt.paragraphs[0].font.size = Pt(12)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].font.color.rgb = color
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER

        if '\n' in desc:
            for line in desc.split('\n'):
                p = bt.add_paragraph()
                p.text = line
                p.font.size = Pt(10)
                p.font.color.rgb = COLORS['text_muted']
                p.alignment = PP_ALIGN.CENTER
        else:
            p = bt.add_paragraph()
            p.text = "\n" + desc
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['text_muted']
            p.alignment = PP_ALIGN.CENTER

        bt.vertical_anchor = MSO_ANCHOR.MIDDLE

    # External services at bottom
    ext_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(8.6), Inches(0.8)
    )
    ext_box.fill.solid()
    ext_box.fill.fore_color.rgb = COLORS['bg_card']
    ext_box.line.color.rgb = COLORS['primary']
    ext_box.line.width = Pt(1.5)

    et = ext_box.text_frame
    et.paragraphs[0].text = "External Services"
    et.paragraphs[0].font.size = Pt(11)
    et.paragraphs[0].font.bold = True
    et.paragraphs[0].font.color.rgb = COLORS['primary']

    p = et.add_paragraph()
    p.text = "AWS APIs (boto3): EC2 | IAM | S3 | VPC | Lambda | RDS    |    AI Provider: Groq/Gemini"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['text_muted']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_graph_modeling_slide(prs):
    """Slide 6: Graph Modeling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "06"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Graph Modeling"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Two column layout
    # Left - Nodes
    nodes_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.3), Inches(4.1), Inches(2.5)
    )
    nodes_box.fill.solid()
    nodes_box.fill.fore_color.rgb = COLORS['bg_card']
    nodes_box.line.color.rgb = COLORS['primary']
    nodes_box.line.width = Pt(2)

    nt = nodes_box.text_frame
    nt.paragraphs[0].text = "Nodes (Cloud Resources)"
    nt.paragraphs[0].font.size = Pt(14)
    nt.paragraphs[0].font.bold = True
    nt.paragraphs[0].font.color.rgb = COLORS['primary']

    node_list = [
        "• IAM: Users, Roles, Policies",
        "• Compute: EC2, Lambda",
        "• Network: VPC, Subnet, Security Group",
        "• Storage: S3, RDS"
    ]
    for node in node_list:
        p = nt.add_paragraph()
        p.text = node
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)

    # Right - Edges
    edges_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.1), Inches(2.5)
    )
    edges_box.fill.solid()
    edges_box.fill.fore_color.rgb = COLORS['bg_card']
    edges_box.line.color.rgb = COLORS['accent']
    edges_box.line.width = Pt(2)

    et = edges_box.text_frame
    et.paragraphs[0].text = "Edges (Relationships)"
    et.paragraphs[0].font.size = Pt(14)
    et.paragraphs[0].font.bold = True
    et.paragraphs[0].font.color.rgb = COLORS['accent']

    edge_list = [
        "• attached_to",
        "• belongs_to",
        "• has_permission",
        "• can_access"
    ]
    for edge in edge_list:
        p = et.add_paragraph()
        p.text = edge
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)

    # Graph visualization - simple and clean
    graph_y = Inches(4.2)
    nodes = [
        (Inches(1.5), graph_y, "IAM User", COLORS['primary']),
        (Inches(3.5), graph_y, "IAM Role", COLORS['primary']),
        (Inches(5.5), graph_y, "EC2", COLORS['success']),
        (Inches(7.5), graph_y, "S3", COLORS['accent']),
    ]

    for x, y, label, color in nodes:
        # Node circle
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(0.9), Inches(0.9)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()

        ntf = node.text_frame
        ntf.paragraphs[0].text = label
        ntf.paragraphs[0].font.size = Pt(9)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = COLORS['bg_dark']
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Connection lines
    connections = [
        (Inches(2.45), graph_y + Inches(0.4), Inches(1), Inches(0.06)),
        (Inches(4.45), graph_y + Inches(0.4), Inches(1), Inches(0.06)),
        (Inches(6.45), graph_y + Inches(0.4), Inches(1), Inches(0.06)),
    ]

    for cx, cy, cw, ch in connections:
        conn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, cy, cw, ch
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = COLORS['text_muted']
        conn.line.fill.background()

    return slide


def add_attack_path_discovery_slide(prs):
    """Slide 7: Attack Path Discovery."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "07"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Attack Path Discovery"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Example attack path - explicit flow diagram
    example_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4)
    )
    tf = example_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Example Attack Path"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # Attack path flow with explicit labels
    path_y = Inches(1.8)
    path_nodes = [
        ("Internet", COLORS['danger']),
        ("Security Group\n(0.0.0.0/0)", COLORS['warning']),
        ("EC2 Instance", COLORS['success']),
        ("IAM Role", COLORS['accent']),
        ("Admin Policy", COLORS['danger']),
    ]

    for i, (label, color) in enumerate(path_nodes):
        x = Inches(0.7 + i * 1.75)

        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, path_y, Inches(1.6), Inches(0.9)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()

        nt = node.text_frame
        if '\n' in label:
            for line in label.split('\n'):
                nt.paragraphs[0].text = line
                nt.paragraphs[0].font.size = Pt(9)
                nt.paragraphs[0].font.bold = True
                nt.paragraphs[0].font.color.rgb = COLORS['bg_dark']
                nt.paragraphs[0].alignment = PP_ALIGN.CENTER
                p = nt.add_paragraph()
        else:
            nt.paragraphs[0].text = label
            nt.paragraphs[0].font.size = Pt(10)
            nt.paragraphs[0].font.bold = True
            nt.paragraphs[0].font.color.rgb = COLORS['bg_dark']
            nt.paragraphs[0].alignment = PP_ALIGN.CENTER
        nt.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Arrow
        if i < len(path_nodes) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                x + Inches(0.65), path_y + Inches(0.95),
                Inches(0.3), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['text_muted']
            arrow.line.fill.background()

    # Detection vs AI clarification
    det_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.8), Inches(4.2), Inches(1.8)
    )
    det_box.fill.solid()
    det_box.fill.fore_color.rgb = COLORS['bg_card']
    det_box.line.color.rgb = COLORS['success']
    det_box.line.width = Pt(2)

    dt = det_box.text_frame
    dt.paragraphs[0].text = "Rule-Based Detection"
    dt.paragraphs[0].font.size = Pt(13)
    dt.paragraphs[0].font.bold = True
    dt.paragraphs[0].font.color.rgb = COLORS['success']

    det_points = [
        "\n• DFS/BFS graph traversal",
        "• Identifies attack paths automatically",
        "• Scores paths by severity"
    ]
    for point in det_points:
        p = dt.add_paragraph()
        p.text = point
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(6)

    # AI box
    ai_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(3.8), Inches(4.2), Inches(1.8)
    )
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = COLORS['bg_card']
    ai_box.line.color.rgb = COLORS['accent']
    ai_box.line.width = Pt(2)

    at = ai_box.text_frame
    at.paragraphs[0].text = "AI Reasoning"
    at.paragraphs[0].font.size = Pt(13)
    at.paragraphs[0].font.bold = True
    at.paragraphs[0].font.color.rgb = COLORS['accent']

    ai_points = [
        "\n• Explains detected paths in natural language",
        "• Prioritizes risks with context",
        "• Generates remediation steps"
    ]
    for point in ai_points:
        p = at.add_paragraph()
        p.text = point
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(6)

    return slide


def add_ai_analysis_slide(prs):
    """Slide 8: AI-Assisted Security Analysis - clarifies detection vs AI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "08"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Assisted Security Analysis"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Clarification text at top
    clar_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.6)
    )
    tf = clar_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Rule-based detection identifies paths. AI explains and prioritizes them."
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_light']
    p.font.name = 'Segoe UI'

    # Two column layout
    # Left - Detection (deterministic)
    det_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(4.2), Inches(2.5)
    )
    det_box.fill.solid()
    det_box.fill.fore_color.rgb = COLORS['bg_card']
    det_box.line.color.rgb = COLORS['success']
    det_box.line.width = Pt(2)

    dt = det_box.text_frame
    dt.paragraphs[0].text = "Rule-Based Detection"
    dt.paragraphs[0].font.size = Pt(15)
    dt.paragraphs[0].font.bold = True
    dt.paragraphs[0].font.color.rgb = COLORS['success']

    det_items = [
        "\n• Graph traversal (DFS/BFS)",
        "• Identifies attack paths",
        "• Computes risk scores",
        "• Deterministic & auditable"
    ]
    for item in det_items:
        p = dt.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)

    # Right - AI Reasoning
    ai_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(2.1), Inches(4.2), Inches(2.5)
    )
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = COLORS['bg_card']
    ai_box.line.color.rgb = COLORS['accent']
    ai_box.line.width = Pt(2)

    at = ai_box.text_frame
    at.paragraphs[0].text = "AI Reasoning Layer"
    at.paragraphs[0].font.size = Pt(15)
    at.paragraphs[0].font.bold = True
    at.paragraphs[0].font.color.rgb = COLORS['accent']

    ai_items = [
        "\n• Explains paths in natural language",
        "• Adds security context (MITRE ATT&CK)",
        "• Prioritizes by business impact",
        "• Generates remediation steps"
    ]
    for item in ai_items:
        p = at.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)

    # Flow arrow
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.95), Inches(3.2), Inches(0.3), Inches(0.4)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS['text_muted']
    arrow.line.fill.background()

    # Bottom - Pipeline visualization
    pipe_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.9), Inches(8.6), Inches(1)
    )
    pipe_box.fill.solid()
    pipe_box.fill.fore_color.rgb = COLORS['bg_card']
    pipe_box.line.color.rgb = COLORS['primary']
    pipe_box.line.width = Pt(1.5)

    pt = pipe_box.text_frame
    pt.paragraphs[0].text = "Pipeline: Graph → Detection → AI → Report"
    pt.paragraphs[0].font.size = Pt(13)
    pt.paragraphs[0].font.bold = True
    pt.paragraphs[0].font.color.rgb = COLORS['primary']
    pt.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = pt.add_paragraph()
    p.text = "\nDetection finds the issues. AI makes them understandable."
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_muted']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

    return slide


def add_evaluation_slide(prs):
    """Slide 9: Evaluation Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "09"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Evaluation Strategy"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # Two column layout
    # Left - Environments
    env_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.3), Inches(4.1), Inches(2.8)
    )
    env_box.fill.solid()
    env_box.fill.fore_color.rgb = COLORS['bg_card']
    env_box.line.color.rgb = COLORS['primary']
    env_box.line.width = Pt(2)

    et = env_box.text_frame
    et.paragraphs[0].text = "Testing Environments"
    et.paragraphs[0].font.size = Pt(15)
    et.paragraphs[0].font.bold = True
    et.paragraphs[0].font.color.rgb = COLORS['primary']

    envs = [
        "\n• CloudGoat (vulnerable AWS)",
        "• Custom misconfigured labs",
        "• OWASP Cloud Security scenarios"
    ]
    for env in envs:
        p = et.add_paragraph()
        p.text = env
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(10)

    # Right - Metrics
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.1), Inches(2.8)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = COLORS['bg_card']
    metrics_box.line.color.rgb = COLORS['success']
    metrics_box.line.width = Pt(2)

    mt = metrics_box.text_frame
    mt.paragraphs[0].text = "Evaluation Metrics"
    mt.paragraphs[0].font.size = Pt(15)
    mt.paragraphs[0].font.bold = True
    mt.paragraphs[0].font.color.rgb = COLORS['success']

    metrics = [
        "\n• Detection Accuracy",
        "• False Positive Rate",
        "• Analysis Time",
        "• AI Explanation Quality"
    ]
    for m in metrics:
        p = mt.add_paragraph()
        p.text = m
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(10)

    # Comparison at bottom
    comp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.4), Inches(8.6), Inches(1.3)
    )
    comp_box.fill.solid()
    comp_box.fill.fore_color.rgb = COLORS['bg_card']
    comp_box.line.color.rgb = COLORS['accent']
    comp_box.line.width = Pt(1.5)

    ct = comp_box.text_frame
    ct.paragraphs[0].text = "Comparison Approach"
    ct.paragraphs[0].font.size = Pt(13)
    ct.paragraphs[0].font.bold = True
    ct.paragraphs[0].font.color.rgb = COLORS['accent']

    p = ct.add_paragraph()
    p.text = "\nRule-based Detection  vs  AI-Assisted Analysis    |    Traditional Scanning  vs  Graph-based Approach"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_muted']
    p.space_before = Pt(6)

    return slide


def add_contributions_slide(prs):
    """Slide 10: Contributions and Impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    add_decorative_elements(slide)

    # Slide number
    slide_num = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.3), Inches(0.5), Inches(0.3)
    )
    tf = slide_num.text_frame
    p = tf.paragraphs[0]
    p.text = "10"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = 'Segoe UI'

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3), Inches(8), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Contributions and Impact"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.font.name = 'Segoe UI'

    # Title underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()

    # 4 contributions in 2x2 grid
    contributions = [
        ("Graph-Based Framework", "Graph-based modeling of cloud infrastructure for attack-path analysis", COLORS['primary'], Inches(0.7), Inches(1.3)),
        ("AI-Assisted Reasoning", "LLM-based explanation and prioritization of detected attack paths", COLORS['accent'], Inches(5.2), Inches(1.3)),
        ("Interactive Visualization", "Real-time graph exploration dashboard for security analysts", COLORS['warning'], Inches(0.7), Inches(3.3)),
        ("Empirical Evaluation", "Systematic testing on CloudGoat and vulnerable cloud labs", COLORS['success'], Inches(5.2), Inches(3.3)),
    ]

    for title, desc, color, x, y in contributions:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.1), Inches(1.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_card']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        ct = card.text_frame
        ct.paragraphs[0].text = title
        ct.paragraphs[0].font.size = Pt(13)
        ct.paragraphs[0].font.bold = True
        ct.paragraphs[0].font.color.rgb = color

        p = ct.add_paragraph()
        p.text = "\n" + desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_muted']
        p.space_before = Pt(6)

    # Future work - full width at bottom
    future_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.3), Inches(8.6), Inches(1.2)
    )
    future_box.fill.solid()
    future_box.fill.fore_color.rgb = COLORS['bg_card']
    future_box.line.color.rgb = COLORS['primary']
    future_box.line.width = Pt(2)

    ft = future_box.text_frame
    ft.paragraphs[0].text = "Future Work"
    ft.paragraphs[0].font.size = Pt(13)
    ft.paragraphs[0].font.bold = True
    ft.paragraphs[0].font.color.rgb = COLORS['primary']

    future_items = [
        "  Multi-cloud support (Azure, GCP)  •  Advanced attack simulation  •  Security automation"
    ]
    for item in future_items:
        p = ft.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(8)

    return slide


def generate_presentation():
    """Generate the 10-slide research presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides = [
        add_title_slide,
        add_problem_statement_slide,
        add_research_question_slide,
        add_proposed_approach_slide,
        add_system_architecture_slide,
        add_graph_modeling_slide,
        add_attack_path_discovery_slide,
        add_ai_analysis_slide,
        add_evaluation_slide,
        add_contributions_slide,
    ]

    for slide_func in slides:
        slide_func(prs)

    output_path = "/home/shaurya/threat-mapper/ppt/cloud_attack_path_analysis_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    generate_presentation()
